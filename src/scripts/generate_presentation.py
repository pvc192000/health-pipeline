#!/usr/bin/env python3
"""Generate visually appealing final CISC 525 presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PROJ = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
OUTPUT = os.path.join(PROJ, "..", "Final_Presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# Colors
DARK_BLUE = RGBColor(25, 42, 86)
ACCENT_BLUE = RGBColor(52, 103, 235)
ACCENT_ORANGE = RGBColor(230, 81, 0)
LIGHT_GRAY = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(40, 40, 40)
SUBTLE_TEXT = RGBColor(120, 120, 120)


def add_bg_rect(slide, color=LIGHT_GRAY):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_accent_bar(slide, y=Inches(0), color=DARK_BLUE, height=Inches(1.2)):
    """Add a colored bar at specified position."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_heading(slide, text, x=Inches(0.8), y=Inches(0.4), size=Pt(32), color=DARK_BLUE, bold=True):
    txBox = slide.shapes.add_textbox(x, y, Inches(11), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    return txBox


def add_body_text(slide, lines, x=Inches(0.8), y=Inches(1.5), width=Inches(11.5), font_size=Pt(14)):
    txBox = slide.shapes.add_textbox(x, y, width, Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("##"):
            p.text = line[2:].strip()
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
            p.space_before = Pt(14)
        elif line.startswith("•"):
            p.text = line
            p.font.size = font_size
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(4)
            p.level = 0
        elif line.startswith("  •"):
            p.text = line.strip()
            p.font.size = Pt(13)
            p.font.color.rgb = SUBTLE_TEXT
            p.space_before = Pt(2)
            p.level = 1
        elif line == "":
            p.text = ""
            p.space_before = Pt(6)
        else:
            p.text = line
            p.font.size = font_size
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(4)
    return txBox


def content_slide(heading, lines):
    slide = prs.slides.add_slide(BLANK)
    add_bg_rect(slide, LIGHT_GRAY)
    # Blue left accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT_BLUE
    strip.line.fill.background()
    add_heading(slide, heading)
    # Thin line under heading
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    add_body_text(slide, lines)
    return slide


def image_slide(heading, image_path, caption=""):
    slide = prs.slides.add_slide(BLANK)
    add_bg_rect(slide, WHITE)
    add_heading(slide, heading, y=Inches(0.2), size=Pt(24))
    if os.path.exists(image_path):
        # Scale image to fit
        from PIL import Image
        img = Image.open(image_path)
        w, h = img.size
        max_w, max_h = Inches(11), Inches(5.0)
        ratio = min(max_w / Emu(w * 914400 // 96), max_h / Emu(h * 914400 // 96))
        img_w = min(max_w, Inches(12))
        left = (prs.slide_width - img_w) // 2
        slide.shapes.add_picture(image_path, Inches(0.7), Inches(1.2), width=min(img_w, Inches(11.5)), height=Inches(5.0))
    if caption:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(11)
        p.font.color.rgb = SUBTLE_TEXT
        p.alignment = PP_ALIGN.CENTER
    return slide


# ============ SLIDES ============

# 1. Title Slide
slide = prs.slides.add_slide(BLANK)
add_bg_rect(slide, DARK_BLUE)
add_accent_bar(slide, y=Inches(5.5), color=ACCENT_BLUE, height=Inches(2))
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "A Cloud-Native Big Data Pipeline for"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Wearable Health Signal Processing\n& Wellness Insights"
p2.font.size = Pt(38)
p2.font.bold = True
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11), Inches(1.2))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "CISC 525 — Big Data Architectures  |  Final Presentation"
p3.font.size = Pt(16)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER
p4 = tf2.add_paragraph()
p4.text = "Param Chokshi  •  Xiukun Hu  |  Spring 2026"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(200, 200, 200)
p4.alignment = PP_ALIGN.CENTER

# 2. Motivation
content_slide("Motivation", [
    "## The Wearable Health Revolution",
    "• Wearable device market hit $120B+ in 2024 — continuous physiological monitoring",
    "• Fitbit Sense captures: HR (every 5s), sleep stages, HRV, SpO2, skin temp, EDA",
    "",
    "## The Gap",
    "• Consumer apps show basic summaries — step count, hours slept",
    "• Deeper signals for stress recovery, illness detection, wellness trends sit unused",
    "• COVID research proved wearables can detect infection before symptoms appear",
    "",
    "## Our Contribution",
    "• End-to-end AWS pipeline: raw sensor data → 4 clinically-informed wellness scores",
    "• Real-time streaming + batch analytics from the same data lake",
    "• Scalable, reproducible, and cost-effective (~$12/month)",
])

# 3. Problem Statement
content_slide("Problem Statement & Objectives", [
    "## Challenge",
    "• 159K+ hourly time-series rows from 71 participants — too large for traditional tools",
    "• Heterogeneous signals (HR, sleep, activity, skin) require unified processing",
    "• No existing open-source pipeline for wearable → wellness scoring at scale",
    "",
    "## Objectives",
    "• Design an AWS-native pipeline supporting both real-time and batch processing",
    "• Predict 4 wellness scores per participant per day:",
    "  • Sleep Quality (0-100) — based on duration, efficiency, deep sleep ratio",
    "  • Stress/Recovery (0-100) — based on HRV, resting HR, stress indicators",
    "  • Activity/Strain (0-100) — based on steps, active minutes, calories",
    "  • Illness Risk (0-100) — based on SpO2, skin temp deviation, HR elevation",
    "",
    "• Evaluate architectural decisions: partitioning, storage format, batch vs. stream",
])

# 4. Data Information
content_slide("Data Information — LifeSnaps Dataset", [
    "## Overview",
    "• Source: Yfantidou et al., Scientific Data, 2022 (CC BY 4.0)",
    "• 71 participants  •  4+ months  •  Fitbit Sense smartwatch",
    "• 159,508 hourly records + 7,410 daily records",
    "",
    "## Signal Types (35+)",
    "• Vitals: Heart Rate, HRV (RMSSD), Resting HR, SpO2",
    "• Sleep: Duration, Efficiency, Deep/Light/REM/Wake ratios",
    "• Activity: Steps, Calories, Distance, Active Zone Minutes",
    "• Skin: Temperature, Electrodermal Activity (EDA)",
    "• Context: Stress score, Breathing rate, VO2max",
    "",
    "## 4 V's of Big Data",
    "• Volume: 9GB+ raw  •  Velocity: 5-second sampling  •  Variety: 35+ types  •  Veracity: 27% no-wear",
])

# 5. Architecture
image_slide(
    "System Architecture — Batch + Stream Paths",
    os.path.join(PROJ, "docs", "figures", "pipeline_architecture_draw_io.png"),
    "Blue = Batch Path (event-driven ETL, 5-min window)  |  Orange = Stream Path (real-time monitoring)  |  Gray = Cross-path data flow"
)

# 6. Tools
content_slide("Tools & Technologies", [
    "## Ingestion",
    "• Kinesis Data Streams + Firehose — real-time ingestion, 60s buffer to S3",
    "",
    "## Processing",
    "• AWS Glue (Spark) — distributed ETL, CSV → Parquet",
    "• EventBridge + Glue Workflows — event-driven orchestration (5-min batch)",
    "",
    "## Storage & Query",
    "• Amazon S3 — 3-tier data lake (Raw / Processed / Curated)",
    "• Athena — serverless SQL on Parquet  •  Glue Data Catalog — metadata",
    "",
    "## ML & Visualization",
    "• SageMaker Processing Jobs — automated retraining + batch inference (triggered by Lambda)",
    "• QuickSight — interactive 3-sheet dashboard (per-user filtering)",
    "• CloudWatch — real-time streaming metrics dashboard",
    "• AWS CDK — Infrastructure as Code, 4 stacks, one-command deploy/teardown",
])

# 7. Experiments
content_slide("Experiments — Pipeline Execution", [
    "## Infrastructure",
    "• CDK deployed 4 stacks in ~5 min (S3, Kinesis, Glue, SageMaker, QuickSight)",
    "",
    "## Batch Pipeline",
    "• Uploaded 22MB CSV → Glue Crawler (63 cols) → ETL → Parquet (88s)",
    "• Athena: 71 participants, avg 8,262 steps/day, 66 bpm resting HR",
    "",
    "## Stream Pipeline",
    "• Producer streamed 159K records at 20/sec → Firehose → S3 (60s buffer)",
    "• CloudWatch: live HR, Steps, Temp, Calories within seconds",
    "",
    "## ML Training",
    "• Lambda triggers SageMaker Processing Job on ETL completion",
    "• 4 GradientBoosting models on 52 features (3,935 training rows)",
    "• Batch inference → wellness scores → S3 curated bucket → QuickSight",
])

# 8. Results - ML
content_slide("Results — Model Performance", [
    "## ML Wellness Scoring (SageMaker Processing Job, GradientBoosting)",
    "",
    "• Sleep Quality:        MAE = 0.85    R² = 0.958",
    "• Stress Recovery:    MAE = 0.85    R² = 0.993",
    "• Activity Strain:      MAE = 0.18    R² = 0.999",
    "• Illness Risk:           MAE = 1.21    R² = 0.985",
    "",
    "## Key Observations",
    "• All models achieve R² > 0.95 — strong predictive power",
    "• Activity Strain is easiest to predict (directly maps to step/activity data)",
    "• Illness Risk has highest MAE (hardest — subtle multi-signal pattern)",
    "",
    "## Top Features",
    "• Sleep: sleep_duration, sleep_deep_ratio, sleep_efficiency",
    "• Stress: rmssd (HRV), resting_hr, stress_score",
    "• Activity: steps, active_minutes, calories",
    "• Illness: spo2, daily_temperature_variation, resting_hr",
])

# 9-11. Dashboard Screenshots
image_slide("Results — Dashboard: Overview",
    os.path.join(PROJ, "docs", "screenshots", "quicksight-overview.png"),
    "Per-participant filtering  •  KPIs (HR, Steps, Calories, Active Min)  •  HR Zones  •  Steps by Time of Day")

image_slide("Results — Dashboard: Activity & Heart",
    os.path.join(PROJ, "docs", "screenshots", "quicksight-activity.png"),
    "Calories trend  •  Active Zone Minutes (Fat Burn / Cardio / Peak)  •  Skin Temperature  •  Distance")

image_slide("Results — Dashboard: Wellness Trends",
    os.path.join(PROJ, "docs", "screenshots", "quicksight-wellness.png"),
    "4 ML Wellness Scores over time  •  Sleep Duration  •  HRV & Stress  •  Resting Heart Rate")

# 12. Cost & Scalability
content_slide("Cost & Scalability", [
    "## Current Cost (71 participants)",
    "• Kinesis: $11/month (1 shard)  •  Glue: $0.13/run  •  Athena: $0.005/query",
    "• Total: ~$12/month (with SageMaker notebook stopped)",
    "",
    "## Scaling Projections",
    "• 10,000 participants → ~$100-150/month",
    "• 1,000,000 participants → ~$3,000-4,000/month",
    "",
    "## Key Architectural Decisions",
    "• Parquet over CSV: 5.5× compression, 5× cheaper queries",
    "• year/month partitioning: optimal for 4-month dataset",
    "• Serverless services: no idle costs, auto-scales",
    "• IaC (CDK): single `cdk destroy --all` removes everything — critical for cost control",
    "",
    "## Storage Comparison",
    "• Raw CSV: 22MB  →  Processed Parquet: 4MB (81% reduction)",
])

# 13. AI Platform Comparison
content_slide("AI Platform Comparison", [
    "## Can this project be replaced by AI/no-code platforms?",
    "",
    "## What CAN be replaced (~30%)",
    "• Dashboards → Power BI, Google Opal, Looker (strong alternatives)",
    "• Simple workflows → n8n, Zapier (webhook triggers)",
    "",
    "## What CANNOT be replaced (~70%)",
    "• Real-time streaming (Kinesis) — platforms are request/response, not stream-oriented",
    "• Distributed Spark ETL (159K rows) — Zapier/n8n process row-by-row",
    "• Custom ML with domain knowledge — AutoML can't encode clinical feature engineering",
    "• Storage optimization (Parquet, partitioning) — platforms abstract away cost control",
    "• Infrastructure as Code — platforms are click-to-configure, no reproducibility",
    "",
    "## Conclusion",
    "• AI platforms excel at visualization, fall short for big data pipelines",
    "• Our hybrid approach: managed AWS services + custom code = best balance",
])

# 14. Live Demo slide
slide = prs.slides.add_slide(BLANK)
add_bg_rect(slide, DARK_BLUE)
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "LIVE DEMO"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = ""
p3 = tf.add_paragraph()
p3.text = "🟠  Stream: Producer → Kinesis → CloudWatch (real-time)"
p3.font.size = Pt(20)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER
p4 = tf.add_paragraph()
p4.text = "🔵  Batch: S3 → Glue → Athena → QuickSight (per-user insights)"
p4.font.size = Pt(20)
p4.font.color.rgb = WHITE
p4.alignment = PP_ALIGN.CENTER

# 15. References
content_slide("References", [
    "## Data",
    "• Yfantidou, S., et al. (2022). LifeSnaps: a 4-month multi-modal dataset",
    "  capturing unobtrusive snapshots of daily life. Scientific Data, 9, 663.",
    "  DOI: 10.5281/zenodo.6826682",
    "",
    "## Research",
    "• Mishra, T., et al. (2020). Pre-symptomatic detection of COVID-19 from",
    "  smartwatch data. Nature Biomedical Engineering, 4, 1208–1220.",
    "• Marz, N. & Warren, J. (2015). Big Data: Principles and Best Practices",
    "  of Scalable Realtime Data Systems. Manning Publications.",
    "",
    "## Tools",
    "• AWS Documentation: Kinesis, Glue, Athena, SageMaker, QuickSight, CDK",
    "• scikit-learn, PySpark, Pandas, python-pptx",
    "",
    "## Project Repository",
    "• GitHub: https://github.com/pvc192000/health-pipeline",
])

# 16. Thank you
slide = prs.slides.add_slide(BLANK)
add_bg_rect(slide, DARK_BLUE)
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
